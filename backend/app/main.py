import json
from fastapi import FastAPI, WebSocket, HTTPException, Depends, UploadFile, File, Form, Request
from fastapi.middleware.cors import CORSMiddleware
from starlette.websockets import WebSocketState
from starlette.websockets import WebSocketDisconnect
from .config import settings
from .stt import STTClient
from .summarizer import Summarizer
from .db import supabase, get_user_lectures
from .auth import (
    UserCreate, UserLogin, register_user, login_user, logout_user, 
    get_current_user, get_authenticated_user_from_header, SupabaseUser, 
    VerifyEmailRequest, AuthResponse, ResendVerificationRequest, 
    resend_verification_email, forgot_password, ForgotPasswordRequest
)
from .user_usages import (
    get_usage, 
    update_usage_plan, 
    update_usage_uploads, 
    update_usage_recordings, 
    get_remaining_uploads_count, 
    get_remaining_recordings_count,
    get_usage_summary,
    reset_user_usage
)
import asyncio
import stripe # Added for Stripe integration
from pydantic import BaseModel # Added for request body model
import docx
import pypdf
from typing import Optional
import base64
import io
import pptx

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# stt_client = STTClient(settings.DEEPGRAM_API_KEY) // Removed global instance
summarizer = Summarizer(settings.OPENAI_API_KEY)

# Auth endpoints
@app.post("/auth/register")
async def register(user_data: UserCreate):
    try:
        return await register_user(user_data)
    except HTTPException as e:
        # Re-raise the HTTPException that was already created in register_user
        raise e
    except Exception as e:
        # This will now only catch truly unexpected errors
        print(f"Unexpected error in /auth/register endpoint: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, 
            detail="An unexpected server error occurred during registration."
        )

@app.post("/auth/login")
async def login(user_data: UserLogin):
    try:
        return await login_user(user_data)
    except Exception as e:
        raise HTTPException(status_code=401, detail=str(e))

@app.post("/auth/logout")
async def logout():
    try:
        return logout_user()
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

class UserProfileResponse(BaseModel):
    id: str
    email: str
    full_name: Optional[str] = None
    subscription_status: Optional[str] = None
    

@app.get("/auth/me", response_model=UserProfileResponse) # Assuming SupabaseUser is the model for the user object
async def read_users_me(current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    # current_user.id is the authenticated user's ID.
    # Now, fetch their profile from your 'profiles' table.
    try:
        profile_response = supabase.table("profiles").select("*").eq("id", current_user.id).single().execute()
        
        profile_data = profile_response.data
        
        # Combine auth info with profile info for the final response
        return UserProfileResponse(
            id=current_user.id,
            email=current_user.email,
            full_name=profile_data.get('full_name'),
            subscription_status=profile_data.get('subscription_status')
            # map other fields here...
        )
    except Exception as e:
        # This can happen if the trigger failed or for a user created before the trigger existed.
        # We can fall back to the auth metadata.
        print(f"Could not find profile for user {current_user.id}, falling back to metadata. Error: {e}")
        return UserProfileResponse(
            id=current_user.id,
            email=current_user.email,
            full_name=current_user.user_metadata.get('full_name') if current_user.user_metadata else None
        )

@app.get("/usage/plan")
async def get_usage_route(current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        return await get_usage(current_user.id)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class UpdateUsageRequest(BaseModel):
    updated_plan: str

@app.put("/usage/update_plan")
async def update_usage_route(
    request: UpdateUsageRequest,
    current_user: SupabaseUser = Depends(get_authenticated_user_from_header)
):
    try:
        return await update_usage_plan(current_user.id, request.updated_plan)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/usage/remaining-uploads")
async def get_remaining_uploads_route(current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        return await get_remaining_uploads_count(current_user.id)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/usage/remaining-recordings")
async def get_remaining_recordings_route(current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        return await get_remaining_recordings_count(current_user.id)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/usage/summary")
async def get_usage_summary_route(current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        return await get_usage_summary(current_user.id)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/lectures")
async def get_lectures(current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        user_id = current_user.id
        lectures_data = get_user_lectures(user_id)
        return lectures_data
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/lectures/{lecture_id}")
async def get_lecture(lecture_id: str, current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        # Query the specific lecture and ensure it belongs to the current user
        response = supabase.table("lectures").select("*").eq("id", lecture_id).eq("user_id", current_user.id).execute()
        
        if not response.data or len(response.data) == 0:
            raise HTTPException(status_code=404, detail="Lecture not found")
            
        return response.data[0]
    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/lectures/{lecture_id}")
async def delete_lecture(lecture_id: str, current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        # First, verify the lecture exists and belongs to the user
        select_response = supabase.table("lectures").select("id").eq("id", lecture_id).eq("user_id", current_user.id).execute()
        
        if not select_response.data:
            raise HTTPException(status_code=404, detail="Lecture not found or you do not have permission to delete it.")

        # If verification passes, proceed with deletion
        delete_response = supabase.table("lectures").delete().eq("id", lecture_id).eq("user_id", current_user.id).execute()

        # The delete operation itself might not return data, so we check for errors if any
        if hasattr(delete_response, 'error') and delete_response.error:
            raise HTTPException(status_code=500, detail=f"Failed to delete lecture: {delete_response.error.message}")

        # Return a success response
        return {"success": True, "message": "Lecture deleted successfully"}

    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred: {str(e)}")

@app.post("/auth/refresh")
async def refresh_token(refresh_token: str):
    try:
        # Use the refresh token to get a new access token
        response = supabase.auth.refresh_session(refresh_token)
        return {
            "access_token": response.session.access_token,
            "refresh_token": response.session.refresh_token
        }
    except Exception as e:
        raise HTTPException(status_code=401, detail=str(e))

@app.post("/auth/verify")
async def verify_email_endpoint(verify_data: VerifyEmailRequest):
    """Verify email address."""
    return await verify_email(verify_data)

@app.post("/auth/resend-verification")
async def resend_verification_endpoint(data: ResendVerificationRequest):
    """Resend verification email."""
    return await resend_verification_email(data)

@app.post("/auth/forgot-password")
async def forgot_password_endpoint(request: ForgotPasswordRequest):
    """Initiate password reset process."""
    return await forgot_password(request)

# Stripe integration
stripe.api_key = settings.STRIPE_SECRET_KEY

class CreateCheckoutSessionRequest(BaseModel):
    plan_type: str # e.g., "plus", "pro", "max"

class UpdateSubscriptionRequest(BaseModel):
    session_id: str

@app.post("/api/v1/stripe/create-checkout-session")
async def create_checkout_session(
    request_data: CreateCheckoutSessionRequest,
    current_user: SupabaseUser = Depends(get_authenticated_user_from_header)
):
    plan_type = request_data.plan_type
    price_id = None

    if plan_type == "plus":
        price_id = settings.STRIPE_PRICE_PLUS
    elif plan_type == "pro":
        price_id = settings.STRIPE_PRICE_PRO
    elif plan_type == "max":
        price_id = settings.STRIPE_PRICE_MAX
    
    if not price_id:
        raise HTTPException(status_code=400, detail="Invalid plan type provided.")

    success_url = f"{settings.FRONTEND_URL}/profile?payment_status=success&session_id={{CHECKOUT_SESSION_ID}}"
    cancel_url = f"{settings.FRONTEND_URL}/profile?payment_status=cancelled"

    try:
        # Check if user already has a stripe_customer_id
        profile_response = supabase.table("profiles").select("stripe_customer_id").eq("id", current_user.id).single().execute()
        customer_id = profile_response.data.get("stripe_customer_id")

        if not customer_id:
            # Create a new Stripe Customer
            customer = stripe.Customer.create(
                email=current_user.email,
                name=current_user.user_metadata.get("full_name"),
                metadata={"supabase_id": current_user.id}
            )
            customer_id = customer.id
            # Save the new customer_id to the user's profile
            supabase.table("profiles").update({"stripe_customer_id": customer_id}).eq("id", current_user.id).execute()

        checkout_session = stripe.checkout.Session.create(
            customer=customer_id,
            line_items=[
                {
                    'price': price_id,
                    'quantity': 1,
                },
            ],
            mode='subscription',
            success_url=success_url,
            cancel_url=cancel_url,
            client_reference_id=current_user.id
        )
        return {"sessionId": checkout_session.id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/stripe/update-subscription")
async def update_subscription_after_payment(
    request_data: UpdateSubscriptionRequest,
    current_user: SupabaseUser = Depends(get_authenticated_user_from_header)
):
    """Update user's subscription status after successful Stripe payment"""
    
    print("Updating subscription endpoint in backend now")
    print("Request data is " + str(request_data))
    print("Current user is " + str(current_user))

    try:
        # Retrieve the checkout session from Stripe to verify payment
        session = stripe.checkout.Session.retrieve(
            request_data.session_id,
            expand=["line_items", "subscription"]
        )

        print("Session is " + str(session))
        
        print("Session Subscription is: " + str(session.subscription))

        if session.payment_status != 'paid':
            raise HTTPException(status_code=400, detail="Payment not completed")
        
        if session.client_reference_id != current_user.id:
            raise HTTPException(status_code=403, detail="Session does not belong to current user")
        
        # Determine the plan type based on the price ID
        price_id = session.line_items.data[0].price.id if session.line_items.data else None
        plan_type = None
        
        if price_id == settings.STRIPE_PRICE_PLUS:
            plan_type = "plus"
        elif price_id == settings.STRIPE_PRICE_PRO:
            plan_type = "pro"
        elif price_id == settings.STRIPE_PRICE_MAX:
            plan_type = "max"
        
        if not plan_type:
            raise HTTPException(status_code=400, detail="Could not determine plan type from session")
        
        # Update the user's subscription status and customer ID in the database
        subscription_id = session.subscription
        customer_id = session.customer

        subscription = session.subscription
        subscription_id = subscription.id

        print("Subscription ID is: " + str(subscription_id))

        print("Subscription Itmes is: " + str(subscription['items']))
        print("Subscription data is: " + str(subscription['items'].data[0]))
        print("Subscription data start date is: " + str(subscription['items'].data[0].current_period_start))
        print("Subscription data end date is: " + str(subscription['items'].data[0].current_period_end))
        
        start_date = subscription['items'].data[0].current_period_start
        end_date = subscription['items'].data[0].current_period_end

        print("Start Date: " + str(start_date))
        print("End Date: " + str(end_date))

        await update_usage_plan(current_user.id, plan_type, stripe_subscription_id=subscription_id, stripe_customer_id=customer_id, start_date=start_date, end_date=end_date)
        
        return {
            "success": True,
            "message": "Subscription updated successfully",
            "plan_type": plan_type
        }
        
    except stripe.error.StripeError as e:
        raise HTTPException(status_code=400, detail=f"Stripe error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/stripe/cancel-subscription")
async def cancel_subscription(current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    """Cancels the user's active Stripe subscription."""
    try:
        # 1. Get the user's stripe_subscription_id from your database
        profile_response = supabase.table("profiles").select("stripe_subscription_id").eq("id", current_user.id).single().execute()
        if not profile_response.data or not profile_response.data.get("stripe_subscription_id"):
            raise HTTPException(status_code=404, detail="No active subscription found to cancel.")
        
        stripe_subscription_id = profile_response.data["stripe_subscription_id"]


        # So now we will set the stripe logic, to instead update the subscription to cancel at the period end time
        stripe.Subscription.modify(
            stripe_subscription_id,
            cancel_at_period_end=True
        )

        return {"success": True, "message": "Your subscription has been successfully canceled."}
    
    except stripe.error.StripeError as e:
        # Handle specific Stripe errors (e.g., subscription already canceled)
        raise HTTPException(status_code=400, detail=f"Stripe error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/stripe/webhook")
async def stripe_webhook(request: Request):
    payload = await request.body()
    sig_header = request.headers.get('stripe-signature')
    event = None

    print(payload)

    try:
        event = stripe.Webhook.construct_event(
            payload, sig_header, settings.STRIPE_WEBHOOK_SECRET
        )
    except ValueError as e:
        # Invalid payload
        raise HTTPException(status_code=400, detail=str(e))
    except stripe.error.SignatureVerificationError as e:
        # Invalid signature
        raise HTTPException(status_code=400, detail=str(e))

    # Handle the event
    if event['type'] == 'invoice.paid':
        invoice = event['data']['object']

        print("Invoice is: " + str(invoice))

        if invoice.billing_reason == 'subscription_cycle':
            print("Subscription Payment Invoice Received")
            stripe_customer_id = invoice.get('customer')

            if stripe_customer_id:
                try:
                    # Find the user with this customer ID
                    user_response = supabase.table("profiles").select("id").eq("stripe_customer_id", stripe_customer_id).single().execute()
                    
                    line_item = invoice['lines']['data'][0]
                    start_date = line_item['period']['start']
                    end_date = line_item['period']['end']

                    print("Start Time is: " + str(start_date))
                    print("End Time is: " + str(end_date))

                    if user_response.data:
                        user_id = user_response.data['id']
                        # Reset their usage
                        await reset_user_usage(user_id, start_date=start_date, end_date=end_date)
                        print(f"Successfully reset usage for user {user_id} via customer ID {stripe_customer_id}")
                    else:
                        print(f"Webhook received for unknown customer: {stripe_customer_id}")
                except Exception as e:
                    print(f"Error processing 'invoice.paid' webhook: {e}")
                    # Return a 500 but don't crash the server, log the error for investigation
                    return {"status": "error", "message": "Internal server error"}
    elif event['type'] == 'customer.subscription.deleted':
        print("Subscription Deleted")
        subscription = event['data']['object']
        stripe_customer_id = subscription.get('customer')

        if stripe_customer_id:
            try:
                user_response = supabase.table("profiles").select("id").eq("stripe_customer_id", stripe_customer_id).single().execute()

                if user_response.data:
                    user_id = user_response.data['id']
                    print(f"Subscription for user {user_id}, has ended. Downgrading to free plan.")

                    await update_usage_plan(
                        user_id=user_id,
                        updated_plan="free",
                        stripe_subscription_id=""
                    )
                else:
                    print(f"Webhook received 'customer.subscription.deleted' for unknown customer: {stripe_customer_id}")
            except Exception as e:
                print(f"Error processing 'customer.subscription.deleted' webhook: {e}")
                return {"status": "error", "message": "Internal server error"}
                
    return {"status": "success"}


@app.websocket("/ws/process-upload")
async def websocket_process_upload(ws: WebSocket):
    await ws.accept()

    try:
        # 1. Authenticate user from token sent by client
        auth_message = await ws.receive_json()
        token = auth_message.get("token")
        if not token:
            await ws.close(code=4001, reason="Missing authentication token")
            return
        
        user_response = supabase.auth.get_user(token)
        if not user_response or not user_response.user:
            await ws.close(code=4001, reason="Invalid authentication token")
            return
        
        user_id = user_response.user.id
        print(f"User {user_id} connected for upload processing.")

        # 2. Receive content (text or file) from the client
        content_message = await ws.receive_json()
        content_type = content_message.get("type")
        
        transcript = ""
        if content_type == "text":
            transcript = content_message.get("data", "")
        elif content_type == "file":
            filename = content_message.get("filename", "")
            file_bytes_str = content_message.get("data", "")
            file_bytes = base64.b64decode(file_bytes_str)

            # Enforce 10MB file size limit
            if len(file_bytes) > 10 * 1024 * 1024:
                await ws.send_text(json.dumps({"error": "File size exceeds the 10MB limit."}))
                return

            file_like_object = io.BytesIO(file_bytes)

            await ws.send_text(json.dumps({
                "processing_status": "Extracting text from your file...",
                "progress": 10
            }))
            
            if filename.endswith(".pdf"):
                reader = pypdf.PdfReader(file_like_object)
                for page in reader.pages:
                    transcript += page.extract_text() or ""
            elif filename.endswith(".docx"):
                doc = docx.Document(file_like_object)
                for para in doc.paragraphs:
                    transcript += para.text + "\n"
            elif filename.endswith(".pptx"):
                prs = pptx.Presentation(file_like_object)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            transcript += shape.text + "\n"
            elif filename.endswith(".txt"):
                transcript = file_like_object.read().decode("utf-8")
            else:
                raise ValueError("Unsupported file type.")
        else:
            raise ValueError("Invalid content type specified.")

        if not transcript.strip():
            raise ValueError("The provided content is empty.")

        # Enforce 15,000 word limit
        words = transcript.split()
        if len(words) > 15000:
            transcript = " ".join(words[:15000])

        # 3. Process the content and send progress updates
        await ws.send_text(json.dumps({
            "processing_status": "Creating your AI-powered summary...",
            "progress": 40
        }))
        summary = await summarizer.summarize(transcript)

        await ws.send_text(json.dumps({
            "processing_status": "Pulling out the important insights...",
            "progress": 60
        }))
        structured_summary_data = summarizer.parse_structured_summary(summary)

        await ws.send_text(json.dumps({
            "processing_status": "Saving your notes securely...",
            "progress": 80
        }))
        
        lecture_data_to_insert = {
            "user_id": user_id,
            "transcript": transcript,
            "summary": summary,
            "lecture_title": structured_summary_data.get("lecture_title", "Uploaded Content"),
            "topic_summary_sentence": structured_summary_data.get("topic_summary_sentence"),
            "key_concepts": structured_summary_data.get("key_concepts"),
            "main_points_covered": structured_summary_data.get("main_points_covered"),
            "conclusion_takeaways": structured_summary_data.get("conclusion_takeaways"),
            "references": structured_summary_data.get("references"),
            "section_summaries": structured_summary_data.get("section_summaries", []),
            "study_questions": structured_summary_data.get("study_questions", []),
            "flashcards": structured_summary_data.get("flashcards", [])
        }

        db_response = supabase.table("lectures").insert(lecture_data_to_insert).execute()
        
        if not db_response.data:
            raise Exception("Failed to save the generated notes.")

        lecture_id = db_response.data[0]['id']

        await ws.send_text(json.dumps({
            "processing_status": "All done! Your notes are ready.",
            "progress": 100,
            "success": True,
            "lecture_id": lecture_id
        }))

        print("Now trying to update the usage uploads for " + user_id)
        await update_usage_uploads(user_id)

    except Exception as e:
        print(f"Error in websocket_process_upload: {str(e)}")
        error_message = str(e)
        if ws.client_state == WebSocketState.CONNECTED:
            await ws.send_text(json.dumps({"error": f"An error occurred: {error_message}"}))
    finally:
        if ws.client_state != WebSocketState.DISCONNECTED:
            await ws.close()
        print("Upload processing WebSocket closed.")

@app.websocket("/ws/transcribe")
async def websocket_transcribe(ws: WebSocket):
    await ws.accept()
    
    stt_client = STTClient(settings.DEEPGRAM_API_KEY) # Create new instance per connection

    print("API_KEY IS: " + settings.DEEPGRAM_API_KEY)

    try:

        token = ws.query_params.get("token")
        if not token:
            raise WebSocketException(code=4001, reason="Missing authentication token")

        user_response = supabase.auth.get_user(token)
        if not user_response or not user_response.user:
            raise WebSocketException(code=4001, reason="Invalid authentication token")

        user_id = user_response.user.id

        print(f"User {user_id} connected for transcription.")
        
        # audio_buffer_for_full_transcript was not used, removed
        audio_buffer = []
        
        first_chunk_received = False # To differentiate timeout before vs. after audio starts
        
        async def audio_gen():
            nonlocal first_chunk_received 
            try:
                while True: # Loop until explicitly broken
                    if ws.client_state != WebSocketState.CONNECTED:
                        print("audio_gen: WebSocket no longer connected. Stopping.")
                        break
                    try:
                        # Timeout to detect when client stops sending audio
                        # or when client sends an empty byte string as an explicit end signal
                        data = await asyncio.wait_for(ws.receive_bytes(), timeout=5.0) # Increased timeout slightly
                        
                        if not data: # Client explicitly sent empty bytes (0-length)
                            print("audio_gen: Received empty bytes from client. Ending audio stream.")
                            break
                        
                        first_chunk_received = True
                        # audio_buffer_for_full_transcript.append(data) # If needed for other purposes
                        yield {'buffer': data}

                    except asyncio.TimeoutError:
                        if first_chunk_received:
                            print("audio_gen: Timeout after receiving audio. Assuming client stopped sending.")
                            break # End of audio from client
                        else:
                            print("audio_gen: Timeout, no audio received yet. Client might be idle or setting up. Continuing to wait.")
                            # Optionally, implement a max idle timeout here if desired
                            continue 
                    except WebSocketDisconnect:
                        print("audio_gen: WebSocket disconnected during receive_bytes.")
                        break # Exit generator
                    except RuntimeError as e: # Handles cases like "Cannot call receive_bytes after connection is closed"
                        if "after connection is closed" in str(e):
                            print(f"audio_gen: WebSocket receive_bytes error after close: {e}")
                        else:
                            print(f"audio_gen: Runtime error during receive_bytes: {e}")
                        break
            finally:
                print("audio_gen: Finished generating audio.")

        final_transcript_segments = []
        last_processed_stt_text = "" # Track the very last text from STT (interim or final)

        # Process audio stream from STTClient
        async for stt_event in stt_client.stream_transcribe(audio_gen()):
            if ws.client_state != WebSocketState.CONNECTED:
                print("WebSocket disconnected while receiving partials. Breaking.")
                break 

            transcript_text = stt_event.get("text", "")
            is_speech_final = stt_event.get("is_speech_final", False)

            if not transcript_text.strip(): # Skip if empty text
                continue

            last_processed_stt_text = transcript_text # Update with every non-empty text

            try:
                await ws.send_text(json.dumps({
                    "text": transcript_text,
                    "is_final_utterance_segment": is_speech_final
                }))
            except Exception as send_err:
                print(f"Error sending partial transcript: {send_err}. Client might have disconnected.")
                break # Stop processing if we can't send to client
            
            if is_speech_final:
                final_transcript_segments.append(transcript_text)
        
        print("Finished iterating stt_client.stream_transcribe.")

        # Construct the final transcript for OpenAI and database
        built_transcript_from_finals = " ".join(final_transcript_segments).strip()
        transcript = built_transcript_from_finals
        
        failed_transcript = False

        if not transcript: # if built_transcript_from_finals is empty
            print("Transcript is empty")
            failed_transcript = True
            if last_processed_stt_text:
                transcript = last_processed_stt_text
            elif first_chunk_received:
                 transcript = "No speech detected in audio"
            else:
                transcript = "Session ended prematurely or no audio sent."

        elif final_transcript_segments and last_processed_stt_text and \
             last_processed_stt_text.startswith(final_transcript_segments[-1]) and \
             len(last_processed_stt_text) > len(final_transcript_segments[-1]):
            # This condition aims to catch if the audio cut off mid-utterance after the last speech_final event.
            # It checks if the last_processed_stt_text (which could be an interim) is a continuation
            # of the last recognized final segment.
            final_transcript_segments[-1] = last_processed_stt_text # Update the last final segment to be more complete
            transcript = " ".join(final_transcript_segments).strip()
        
        print("Full Transcript for OpenAI: " + transcript)

        summary = "Summary could not be generated for this session."

        try:
            # Send processing status
            await ws.send_text(json.dumps({
                "processing_status": "Getting everything ready for analysis...",
                "progress": 20
            }))

            # Generate summary
            await ws.send_text(json.dumps({
                "processing_status": "Creating your AI-powered summary...",
                "progress": 40
            }))

            if not failed_transcript:
                summary = await summarizer.summarize(transcript)
                print("\nChatGPT Full Summary: " + summary)

                await ws.send_text(json.dumps({
                    "processing_status": "Pulling out the important insights...",
                    "progress": 60
                }))
                structured_summary_data = summarizer.parse_structured_summary(summary)
                print("\nParsed Structured Summary Data:\n", structured_summary_data)

            # Store in DB
            await ws.send_text(json.dumps({
                "processing_status": "Saving your notes securely...",
                "progress": 80
            }))
            
            # Define conditions for not storing, e.g., very short or error messages
            skippable_transcripts = [
                "No speech detected in audio", 
                "Session ended prematurely or no audio sent."
            ]
            # Also consider not storing if summary indicates an error.
            if transcript.strip() and transcript not in skippable_transcripts and "Error generating summary" not in summary:
                lecture_data_to_insert = {
                    "user_id": user_id,
                    "transcript": transcript,
                    "summary": summary, 
                    "lecture_title": structured_summary_data.get("lecture_title"),
                    "topic_summary_sentence": structured_summary_data.get("topic_summary_sentence"),
                    "key_concepts": structured_summary_data.get("key_concepts", []),
                    "main_points_covered": structured_summary_data.get("main_points_covered", []),
                    "conclusion_takeaways": structured_summary_data.get("conclusion_takeaways"),
                    "references": structured_summary_data.get("references", []),
                    "section_summaries": structured_summary_data.get("section_summaries", []),
                    "study_questions": structured_summary_data.get("study_questions", []),
                    "flashcards": structured_summary_data.get("flashcards", [])
                }

                db_response = supabase.table("lectures").insert(lecture_data_to_insert).execute()
                print(f"Data insertion response: {db_response}")
                
                # Get the lecture ID from the response
                lecture_id = db_response.data[0]['id'] if db_response.data else None

            else:
                print(f"Skipping DB insert for transcript: '{transcript}' or due to summary error.")
                lecture_id = None

            await ws.send_text(json.dumps({
                "processing_status": "Putting on the final touches...",
                "progress": 95
            }))

            # Ensure summary and full transcript are sent if connection is still open
            if ws.client_state == WebSocketState.CONNECTED:
                try:
                    await ws.send_text(json.dumps({
                        "summary": summary,
                        "transcript": transcript,
                        "lecture_id": lecture_id,
                        "processing_status": "All done! Your notes are ready.",
                        "progress": 100
                    }))
                    print("Summary and final transcript sent successfully.")

                    print("Now trying to increment the usage recordings for " + user_id)
                    if not failed_transcript:
                        await update_usage_recordings(user_id)
                except Exception as e:
                    print(f"Failed to send summary/final transcript: {str(e)}")
            else:
                print("Client disconnected, couldn't send summary/final transcript.")

        except Exception as e:
            print(f"Error in summary generation or database storage: {str(e)}")
            # Try to send an error message to the client if still connected
            if ws.client_state == WebSocketState.CONNECTED:
                try:
                    await ws.send_text(json.dumps({
                        "error": "Oops! Something went wrong while saving or summarizing.",
                        "transcript": transcript # Send transcript even if summary fails
                    }))
                except Exception as send_err:
                    print(f"Failed to send error message to client: {send_err}")
            
    except WebSocketDisconnect:
        print(f"Client disconnected from FastAPI WebSocket: {ws.client_state}, {ws.application_state}")
    except Exception as e:
        print(f"Unhandled error in websocket_transcribe: {str(e)}")
        if ws.client_state == WebSocketState.CONNECTED:
            try:
                await ws.send_text(json.dumps({"error": f"An unexpected server error occurred: {str(e)}"}))
            except Exception as send_err:
                print(f"Error sending unhandled error to client: {send_err}")
    finally:
        print("FastAPI WebSocket handler: Attempting to clean up WebSocket.")
        if ws.client_state != WebSocketState.DISCONNECTED:
            try:
                await ws.close(code=1000) 
                print("FastAPI WebSocket closed in finally block.")
            except RuntimeError as e:
                # Handle cases where close() is called on an already closing/closed socket
                print(f"RuntimeError during WebSocket close: {e}. Connection state: {ws.client_state}")
            except Exception as e:
                print(f"Unexpected error during WebSocket close: {e}")
        else:
            print("FastAPI WebSocket already disconnected.")

@app.get("/user/stats")
async def get_user_stats(current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        # Get user's lectures
        response = supabase.table("lectures").select("*").eq("user_id", current_user.id).execute()
        lectures = response.data if response.data else []
        
        # Calculate total lectures and total minutes (assuming average of 5 minutes per lecture for now)
        total_lectures = len(lectures)
        total_minutes = total_lectures * 5  # This is a placeholder. You might want to store actual duration in your lectures table
        
        return {
            "total_lectures": total_lectures,
            "total_minutes": total_minutes
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/lectures/{lecture_id}/flashcards")
async def create_flashcards(lecture_id: str, current_user: SupabaseUser = Depends(get_authenticated_user_from_header)):
    try:
        # 1. Fetch the lecture to ensure it exists and belongs to the user
        response = supabase.table("lectures").select("transcript").eq("id", lecture_id).eq("user_id", current_user.id).execute()
        
        if not response.data:
            raise HTTPException(status_code=404, detail="Lecture not found")
        
        transcript = response.data[0].get("transcript")
        if not transcript or not transcript.strip():
            raise HTTPException(status_code=400, detail="Lecture has no content to generate flashcards from.")

        # 2. Generate flashcards using the summarizer
        flashcards = await summarizer.generate_flashcards(transcript)

        if not flashcards:
            raise HTTPException(status_code=500, detail="Failed to generate flashcards.")

        return flashcards
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred: {str(e)}")